"""
UNIVERSAL DATASET DOWNLOADER v5 - UNIVERSAL FORMAT MODE
==========================================================

Purpose:
    Download datasets from common sources and save them as:
        1. Excel (.xlsx)
        2. CSV (.csv)
        3. Parquet (.parquet)

Supported:
    - Hugging Face dataset ID
        qiaojin/PubMedQA
    - Hugging Face dataset URL
        https://huggingface.co/datasets/qiaojin/PubMedQA
    - Hugging Face direct file URL
    - GitHub normal/blob URL
        https://github.com/user/repo/blob/main/data/file.tsv
    - GitHub raw URL
        https://raw.githubusercontent.com/user/repo/main/data/file.tsv
    - Direct CSV / TSV / JSON / JSONL / Parquet / Excel / ZIP URLs
    - Local CSV/TSV/JSON/JSONL/Parquet/Excel/ZIP/TAR/TAR.GZ files
    - GitHub Markdown/README documentation that points to a dataset download script or archive
    - Markdown tables (.md)

IMPORTANT:
    GitHub "blob" pages are HTML webpages, NOT the dataset itself.
    This version automatically converts GitHub blob URLs to raw URLs
    BEFORE downloading and also checks downloaded content so an HTML
    GitHub page is not accidentally saved as dataset data.

Install once (recommended full support):
    pip install pandas openpyxl requests datasets huggingface_hub pyarrow beautifulsoup4 lxml python-docx python-pptx pypdf pyxlsb odfpy scipy numpy pyyaml tables fastavro dbfread

Run:
    python universal_dataset_to_excel.py
"""

import os
import re
import json
import gzip
import bz2
import lzma
import io
import time
import hashlib
import mimetypes
import sqlite3
import struct
import shutil
import zipfile
import tarfile
import tempfile
from pathlib import Path
from urllib.parse import urlparse, unquote

import pandas as pd
import requests


# ============================================================
# SETTINGS
# ============================================================

OUTPUT_DIR = Path("downloaded_datasets")
CACHE_DIR = Path(".dataset_cache")

OUTPUT_DIR.mkdir(exist_ok=True)
CACHE_DIR.mkdir(exist_ok=True)

EXCEL_MAX_ROWS = 1_000_000
REQUEST_TIMEOUT = 60
MAX_RETRIES = 6

USER_AGENT = (
    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
    "AppleWebKit/537.36 (KHTML, like Gecko) "
    "Chrome/151.0 Safari/537.36"
)

SAVE_EXTRA = True

session = requests.Session()
session.headers.update({
    "User-Agent": USER_AGENT,
    "Accept": "*/*",
    "Accept-Language": "en-US,en;q=0.9",
})


# ============================================================
# BASIC HELPERS
# ============================================================

def safe_name(value):
    value = str(value).strip()
    value = re.sub(r'[<>:"/\\|?*]+', "_", value)
    value = re.sub(r"\s+", "_", value)
    return value[:180] or "dataset"


def unique_path(path):
    if not path.exists():
        return path

    i = 2
    while True:
        candidate = path.with_name(
            f"{path.stem}_{i}{path.suffix}"
        )
        if not candidate.exists():
            return candidate
        i += 1


def is_url(value):
    try:
        p = urlparse(value)
        return p.scheme.lower() in ("http", "https") and bool(p.netloc)
    except Exception:
        return False


def is_github_url(url):
    try:
        return urlparse(url).netloc.lower() in {
            "github.com",
            "www.github.com",
            "raw.githubusercontent.com",
        }
    except Exception:
        return False


def github_to_raw(url):
    """
    Convert GitHub webpage/blob URL to raw.githubusercontent.com URL.

    Example:
        https://github.com/user/repo/blob/main/data/file.tsv

    becomes:
        https://raw.githubusercontent.com/user/repo/main/data/file.tsv
    """

    url = url.strip()

    parsed = urlparse(url)
    host = parsed.netloc.lower()
    path = parsed.path.strip("/")

    # Already raw
    if host == "raw.githubusercontent.com":
        return url

    if host not in ("github.com", "www.github.com"):
        return url

    # /user/repo/blob/branch/path/file
    parts = path.split("/")

    if len(parts) >= 5 and parts[2] in ("blob", "raw"):
        user = parts[0]
        repo = parts[1]
        mode = parts[2]
        remaining = parts[3:]

        # Correct raw URL does not contain "blob" or "raw".
        raw_path = "/".join([user, repo] + remaining)

        return f"https://raw.githubusercontent.com/{raw_path}"

    # Also handle unusual URL strings containing /blob/
    if "/blob/" in path:
        new_path = path.replace("/blob/", "/", 1)
        return f"https://raw.githubusercontent.com/{new_path}"

    return url


def hf_dataset_url_to_id(url):
    parsed = urlparse(url)
    path = parsed.path.strip("/")

    parts = path.split("/")

    if "datasets" not in parts:
        return None

    i = parts.index("datasets")

    if len(parts) <= i + 2:
        return None

    return f"{parts[i+1]}/{parts[i+2]}"


def looks_like_hf_id(value):
    if is_url(value):
        return False

    return bool(
        re.match(
            r"^[A-Za-z0-9_.-]+/[A-Za-z0-9_.-]+$",
            value.strip()
        )
    )


def filename_from_url(url):
    name = Path(
        unquote(urlparse(url).path)
    ).name

    return name or "dataset"


def dataset_name_from_url(url):
    name = filename_from_url(url)

    for ext in [
        ".jsonl.gz",
        ".tar.gz",
        ".jsonl",
        ".parquet",
        ".xlsx",
        ".xls",
        ".csv",
        ".tsv",
        ".txt",
        ".json",
        ".zip",
        ".tar",
        ".tgz",
        ".md",
    ]:
        if name.lower().endswith(ext):
            name = name[:-len(ext)]
            break

    return safe_name(name)


# ============================================================
# HTTP DOWNLOAD
# ============================================================

def download_http(url):
    """
    Download with retry.

    Special handling:
        429 = rate limit
        500/502/503/504 = temporary server problem
    """

    last_error = None

    for attempt in range(1, MAX_RETRIES + 1):

        try:
            print(f"\nDownloading URL:")
            print(url)

            response = session.get(
                url,
                timeout=REQUEST_TIMEOUT,
                stream=True,
                allow_redirects=True,
            )

            if response.status_code == 429:
                retry_after = response.headers.get("Retry-After")

                if retry_after and retry_after.isdigit():
                    wait = int(retry_after)
                else:
                    wait = min(90, attempt * 5)

                print(
                    f"HTTP 429 - rate limited. "
                    f"Waiting {wait} seconds..."
                )

                response.close()
                time.sleep(wait)
                continue

            if response.status_code in (500, 502, 503, 504):
                wait = min(60, attempt * 4)

                print(
                    f"HTTP {response.status_code}. "
                    f"Retrying in {wait} seconds..."
                )

                response.close()
                time.sleep(wait)
                continue

            response.raise_for_status()

            content_type = response.headers.get(
                "Content-Type", ""
            ).lower()

            final_url = response.url

            print("Final URL:", final_url)
            print("Content-Type:", content_type)

            filename = filename_from_url(final_url)

            disposition = response.headers.get(
                "Content-Disposition", ""
            )

            match = re.search(
                r'filename="?([^";]+)"?',
                disposition,
                re.I
            )

            if match:
                filename = unquote(match.group(1))

            filename = safe_name(filename)

            # If GitHub raw URL has no useful extension,
            # keep a temporary .bin file and later detect content.
            target = unique_path(
                CACHE_DIR / filename
            )

            total = int(
                response.headers.get(
                    "Content-Length", 0
                ) or 0
            )

            downloaded = 0

            with open(target, "wb") as f:

                for chunk in response.iter_content(
                    chunk_size=1024 * 1024
                ):

                    if not chunk:
                        continue

                    f.write(chunk)
                    downloaded += len(chunk)

                    if total:
                        pct = downloaded * 100 / total
                        print(
                            f"\rDownloaded "
                            f"{downloaded / 1024 / 1024:.2f} MB "
                            f"({pct:.1f}%)",
                            end=""
                        )
                    else:
                        print(
                            f"\rDownloaded "
                            f"{downloaded / 1024 / 1024:.2f} MB",
                            end=""
                        )

            response.close()

            print("\nSaved:", target)

            # IMPORTANT:
            # Check whether GitHub accidentally returned HTML.
            if is_html_file(target):

                print(
                    "\nWARNING: Server returned HTML instead "
                    "of dataset data."
                )

                # If original was GitHub blob URL, force raw.
                if is_github_url(url):

                    raw_url = github_to_raw(url)

                    if raw_url != url:

                        print(
                            "GitHub webpage detected."
                        )
                        print(
                            "Trying RAW URL automatically:"
                        )
                        print(raw_url)

                        return download_http(
                            raw_url
                        )

                raise RuntimeError(
                    "The URL returned an HTML webpage, "
                    "not a dataset file."
                )

            return target

        except requests.RequestException as exc:

            last_error = exc

            wait = min(60, attempt * 4)

            print(
                f"\nConnection error: {exc}"
            )

            if attempt < MAX_RETRIES:
                print(
                    f"Retrying in {wait} seconds..."
                )
                time.sleep(wait)

    raise RuntimeError(
        f"Download failed after {MAX_RETRIES} attempts.\n"
        f"URL: {url}\n"
        f"Last error: {last_error}"
    )


# ============================================================
# HTML DETECTION
# ============================================================

def is_html_file(path):
    """
    Detect GitHub/login/error HTML before treating it as TSV/CSV.
    """

    try:
        with open(path, "rb") as f:
            sample = f.read(100_000)

        text = sample.decode(
            "utf-8",
            errors="ignore"
        ).lstrip().lower()

        # Strong HTML signatures
        if text.startswith("<!doctype html"):
            return True

        if text.startswith("<html"):
            return True

        if "<html" in text[:5000]:
            return True

        if "<head" in text[:5000] and "<body" in text[:10000]:
            return True

        # GitHub page markers
        github_markers = [
            "github.com",
            "sign in to github",
            "skip to content",
            "githubusercontent",
        ]

        if "<title>" in text:
            if any(
                marker in text
                for marker in github_markers
            ):
                return True

    except Exception:
        pass

    return False



# ============================================================
# ARCHIVE / MARKDOWN HELPERS
# ============================================================

ARCHIVE_SUFFIXES = (
    ".zip",
    ".tar",
    ".tar.gz",
    ".tgz",
    ".tar.bz2",
    ".tbz2",
    ".tar.xz",
    ".txz",
)

DATA_SUFFIXES = (
    ".csv", ".tsv", ".txt", ".log", ".data", ".dat",
    ".json", ".jsonl", ".ndjson", ".jsonl.gz",
    ".yaml", ".yml", ".xml", ".md", ".markdown",
    ".html", ".htm", ".parquet", ".feather", ".orc",
    ".xlsx", ".xls", ".xlsb", ".ods",
    ".h5", ".hdf5", ".hdf",
    ".dta", ".sav", ".arff", ".mat", ".npy", ".npz",
    ".db", ".sqlite", ".sqlite3", ".db3",
    ".zip", ".tar", ".tar.gz", ".tgz", ".tar.bz2",
    ".tbz2", ".tar.xz", ".txz",
)

def is_archive(path):
    """
    Detect archives by CONTENT as well as extension.
    This is important for URLs such as /download/file with no extension.
    """
    name = path.name.lower()

    if name.endswith(ARCHIVE_SUFFIXES):
        return True

    try:
        if zipfile.is_zipfile(path):
            return True
    except Exception:
        pass

    try:
        if tarfile.is_tarfile(path):
            return True
    except Exception:
        pass

    return False


def is_office_zip(path):
    """Distinguish DOCX/PPTX/XLSX containers from ordinary ZIP datasets."""
    try:
        if not zipfile.is_zipfile(path):
            return False
        with zipfile.ZipFile(path, "r") as z:
            names = set(z.namelist())
        return (
            "[Content_Types].xml" in names
            and (
                any(n.startswith("word/") for n in names)
                or any(n.startswith("ppt/") for n in names)
                or any(n.startswith("xl/") for n in names)
            )
        )
    except Exception:
        return False


def markdown_table_to_dataframe(lines):
    """
    Parse one GitHub/CommonMark-style Markdown table.
    Returns a DataFrame or None.
    """
    cleaned = [line.strip() for line in lines if line.strip()]
    if len(cleaned) < 2:
        return None

    for i in range(len(cleaned) - 1):
        header = cleaned[i]
        separator = cleaned[i + 1]

        if "|" not in header or "|" not in separator:
            continue

        sep_cells = [
            cell.strip()
            for cell in separator.strip("|").split("|")
        ]

        if not sep_cells or not all(
            re.fullmatch(r":?-{3,}:?", cell)
            for cell in sep_cells
        ):
            continue

        headers = [
            cell.strip()
            for cell in header.strip("|").split("|")
        ]

        if len(headers) != len(sep_cells):
            continue

        rows = []
        j = i + 2

        while j < len(cleaned) and "|" in cleaned[j]:
            cells = [
                cell.strip()
                for cell in cleaned[j].strip("|").split("|")
            ]

            if len(cells) != len(headers):
                break

            rows.append(cells)
            j += 1

        return pd.DataFrame(rows, columns=headers)

    return None


def read_markdown(path):
    """
    Read Markdown documentation.

    If the Markdown contains a real table, return it as a DataFrame.
    Otherwise raise a clear error instead of pretending the document
    itself is a dataset.
    """
    text = read_text_sample(path, size=2_000_000)

    # Try every possible table starting point.
    lines = text.splitlines()
    df = markdown_table_to_dataframe(lines)

    if df is not None:
        return df

    raise ValueError(
        "Markdown file is documentation, not a tabular dataset. "
        "If this is a GitHub data.md/README URL, the downloader can "
        "automatically follow a dataset download script when possible. "
        "For a local Markdown file, provide the actual CSV/JSONL/Parquet/"
        "Excel/archive dataset file."
    )


def extract_urls_from_text(text):
    """
    Extract HTTP(S) URLs from Markdown, shell scripts, and plain text.
    Only return URLs that look like downloadable data/archive files.
    """
    urls = re.findall(r'https?://[^\s<>"\'`)]+', text)
    result = []

    for url in urls:
        url = url.rstrip(".,;")
        lower = url.lower()

        # Strip common shell/markdown wrappers.
        url = url.rstrip('\'"')

        if any(
            ext in lower
            for ext in DATA_SUFFIXES
        ):
            result.append(url)

    # Preserve order and remove duplicates.
    return list(dict.fromkeys(result))


def github_raw_root_from_url(url):
    """
    Convert:
      https://github.com/user/repo/blob/main/doc/data.md
    or:
      https://raw.githubusercontent.com/user/repo/main/doc/data.md
    into:
      https://raw.githubusercontent.com/user/repo/main
    """
    raw = github_to_raw(url)
    parsed = urlparse(raw)
    parts = parsed.path.strip("/").split("/")

    if len(parts) < 4:
        return None

    user, repo, branch = parts[:3]
    return f"https://raw.githubusercontent.com/{user}/{repo}/{branch}"


def download_small_text(url):
    """
    Download a small documentation/script file into CACHE_DIR and return
    its decoded text. Uses the same session and retry behavior as the
    normal downloader, but avoids treating .md/.sh as datasets.
    """
    last_error = None

    for attempt in range(1, MAX_RETRIES + 1):
        try:
            response = session.get(
                url,
                timeout=REQUEST_TIMEOUT,
                allow_redirects=True,
            )
            response.raise_for_status()

            content = response.content
            response.close()

            return content.decode(
                "utf-8-sig",
                errors="replace",
            )

        except requests.RequestException as exc:
            last_error = exc
            wait = min(60, attempt * 4)

            if attempt < MAX_RETRIES:
                print(
                    f"\nDocumentation download error: {exc}"
                    f"\nRetrying in {wait} seconds..."
                )
                time.sleep(wait)

    raise RuntimeError(
        f"Could not read documentation URL after "
        f"{MAX_RETRIES} attempts: {url}\n"
        f"Last error: {last_error}"
    )


def find_github_dataset_urls(document_url):
    """
    For a GitHub Markdown/README file:
      1. Look for direct downloadable dataset/archive URLs.
      2. If none are present, look for the repository's script/get_data.sh
         and extract its dataset/archive URLs.

    This specifically fixes cases where users paste a GitHub data.md file
    that describes the dataset but is not itself the dataset.
    """
    raw_url = github_to_raw(document_url)

    print("\nReading GitHub documentation:")
    print(raw_url)

    text = download_small_text(raw_url)

    urls = extract_urls_from_text(text)

    # For README/data.md files, try the common repository download script
    # if the Markdown itself does not expose a direct dataset URL.
    if not urls:
        root = github_raw_root_from_url(document_url)

        if root:
            script_candidates = [
                f"{root}/script/get_data.sh",
                f"{root}/scripts/get_data.sh",
                f"{root}/get_data.sh",
                f"{root}/download_data.sh",
            ]

            for script_url in script_candidates:
                try:
                    script_text = download_small_text(script_url)
                    script_urls = extract_urls_from_text(script_text)

                    if script_urls:
                        print(
                            "\nDataset download script detected:"
                        )
                        print(script_url)
                        urls.extend(script_urls)
                        break

                except Exception:
                    continue

    return list(dict.fromkeys(urls))



# ============================================================
# FILE READERS
# ============================================================

def read_text_sample(path, size=30_000):

    with open(path, "rb") as f:
        raw = f.read(size)

    for encoding in (
        "utf-8-sig",
        "utf-8",
        "utf-16",
        "cp1252",
        "latin-1",
    ):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue

    return raw.decode(
        "utf-8",
        errors="replace"
    )


def read_csv_tsv(path):

    suffix = path.suffix.lower()

    if suffix == ".tsv":
        return pd.read_csv(
            path,
            sep="\t",
            encoding="utf-8-sig"
        )

    # Automatic delimiter detection
    try:
        return pd.read_csv(
            path,
            sep=None,
            engine="python",
            encoding="utf-8-sig"
        )
    except Exception:
        pass

    # Fallback delimiters
    for sep in [",", "\t", ";", "|"]:

        try:
            return pd.read_csv(
                path,
                sep=sep,
                encoding="utf-8-sig"
            )
        except Exception:
            continue

    raise ValueError(
        f"Could not read CSV/TSV file: {path}"
    )


def read_json(path):

    name = path.name.lower()

    if name.endswith(".jsonl.gz"):

        with gzip.open(
            path,
            "rt",
            encoding="utf-8-sig"
        ) as f:
            records = [
                json.loads(line)
                for line in f
                if line.strip()
            ]

        return pd.json_normalize(records)

    text = path.read_text(
        encoding="utf-8-sig"
    )

    # JSONL
    if name.endswith(".jsonl"):

        records = [
            json.loads(line)
            for line in text.splitlines()
            if line.strip()
        ]

        return pd.json_normalize(records)

    obj = json.loads(text)

    if isinstance(obj, list):
        return pd.json_normalize(obj)

    if isinstance(obj, dict):

        for key in [
            "data",
            "records",
            "items",
            "results",
            "dataset",
        ]:

            if (
                key in obj
                and isinstance(obj[key], list)
            ):
                return pd.json_normalize(
                    obj[key]
                )

        return pd.json_normalize(obj)

    raise ValueError(
        "Unsupported JSON structure."
    )



# ============================================================
# UNIVERSAL CONTENT DETECTION / READERS
# ============================================================

def file_signature(path):
    """Return a short binary signature for format detection."""
    try:
        with open(path, "rb") as f:
            return f.read(32)
    except Exception:
        return b""


def looks_like_sqlite(path):
    return file_signature(path).startswith(b"SQLite format 3\x00")


def looks_like_pdf(path):
    return file_signature(path).startswith(b"%PDF-")


def looks_like_gzip(path):
    return file_signature(path).startswith(b"\x1f\x8b")


def looks_like_bzip2(path):
    return file_signature(path).startswith(b"BZh")


def looks_like_xz(path):
    return file_signature(path).startswith(b"\xfd7zXZ\x00")


def looks_like_text(path, sample_size=200_000):
    """
    Content-based text detector. Extension is deliberately ignored.
    """
    try:
        with open(path, "rb") as f:
            raw = f.read(sample_size)
        if not raw:
            return True
        if b"\x00" in raw:
            return False

        for enc in ("utf-8", "utf-16", "cp1252", "latin-1"):
            try:
                decoded = raw.decode(enc)
                bad = sum(1 for ch in decoded if ord(ch) < 32 and ch not in "\r\n\t")
                return bad / max(len(decoded), 1) < 0.01
            except UnicodeDecodeError:
                continue
    except Exception:
        pass
    return False


def read_delimited_universal(path):
    """
    Try delimiter inference for unknown-extension tabular text.
    """
    encodings = ("utf-8-sig", "utf-8", "utf-16", "cp1252", "latin-1")
    errors = []

    for encoding in encodings:
        try:
            # First try automatic delimiter detection.
            df = pd.read_csv(
                path,
                sep=None,
                engine="python",
                encoding=encoding,
                on_bad_lines="warn",
            )
            if df.shape[1] > 1 or len(df) > 0:
                return df
        except Exception as exc:
            errors.append(exc)

        # Common delimiters as a fallback.
        for sep in [",", "\t", ";", "|", ":", " "]:
            try:
                df = pd.read_csv(
                    path,
                    sep=sep,
                    engine="python",
                    encoding=encoding,
                    on_bad_lines="warn",
                )
                if df.shape[1] > 1:
                    return df
            except Exception:
                pass

    raise ValueError("Not a readable delimited text table.")


def read_plain_text_as_dataframe(path):
    """
    Last text fallback. Every line becomes a row, so even an unknown
    extension such as .log/.py/.sql/.xyz can still be exported to Excel.
    """
    text = read_text_sample(path, size=5_000_000)
    lines = text.splitlines()

    if not lines:
        return pd.DataFrame({"text": []})

    return pd.DataFrame(
        {
            "line_number": range(1, len(lines) + 1),
            "text": lines,
        }
    )


def read_yaml(path):
    try:
        import yaml
    except ImportError:
        raise RuntimeError("YAML support requires: pip install pyyaml")

    text = path.read_text(encoding="utf-8-sig", errors="replace")
    obj = yaml.safe_load(text)

    if isinstance(obj, list):
        return pd.json_normalize(obj)

    if isinstance(obj, dict):
        for key in ("data", "records", "items", "results", "dataset"):
            if isinstance(obj.get(key), list):
                return pd.json_normalize(obj[key])
        return pd.json_normalize(obj)

    return pd.DataFrame({"value": [obj]})


def xml_element_to_record(element):
    """
    Convert one XML element into a flat dictionary.
    Child elements become columns; attributes get @ prefixes.
    """
    record = {}
    for key, value in element.attrib.items():
        record[f"@{key}"] = value

    for child in list(element):
        key = child.tag.split("}")[-1]
        value = "".join(child.itertext()).strip()
        if key in record:
            # Preserve repeated fields rather than silently dropping them.
            old = record[key]
            if not isinstance(old, list):
                old = [old]
            old.append(value)
            record[key] = old
        else:
            record[key] = value

    if not record:
        record["text"] = (element.text or "").strip()

    return record


def read_xml(path):
    import xml.etree.ElementTree as ET

    root = ET.parse(path).getroot()

    # Prefer repeated child elements as records.
    children = list(root)
    if children:
        records = [xml_element_to_record(x) for x in children]
        if any(record for record in records):
            return pd.json_normalize(records)

    return pd.json_normalize([xml_element_to_record(root)])


def read_sqlite(path):
    result = []
    con = sqlite3.connect(str(path))
    try:
        tables = pd.read_sql_query(
            "SELECT name FROM sqlite_master "
            "WHERE type='table' AND name NOT LIKE 'sqlite_%' "
            "ORDER BY name",
            con,
        )

        for table in tables["name"].tolist():
            # Quote table name safely for SQLite identifiers.
            safe_table = str(table).replace('"', '""')
            df = pd.read_sql_query(
                f'SELECT * FROM "{safe_table}"',
                con,
            )
            result.append((safe_name(table), df))

        if not result:
            result.append(
                (
                    "database_info",
                    pd.DataFrame(
                        {
                            "message": [
                                "SQLite database contains no user tables."
                            ]
                        }
                    ),
                )
            )

        return result
    finally:
        con.close()


def read_excel_universal(path):
    """
    Excel-family reader with XLSX/XLS/XLSB/ODS support.
    """
    suffix = path.suffix.lower()

    kwargs = {}
    if suffix == ".xlsb":
        kwargs["engine"] = "pyxlsb"
    elif suffix == ".ods":
        kwargs["engine"] = "odf"

    excel = pd.ExcelFile(path, **kwargs)
    result = []

    for sheet in excel.sheet_names:
        df = pd.read_excel(
            path,
            sheet_name=sheet,
            **kwargs,
        )
        result.append((safe_name(sheet), df))

    return result


def read_docx(path):
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("DOCX support requires: pip install python-docx")

    doc = Document(path)
    result = []

    paragraphs = [
        p.text for p in doc.paragraphs
        if p.text.strip()
    ]

    if paragraphs:
        result.append(
            (
                "paragraphs",
                pd.DataFrame(
                    {
                        "paragraph_number": range(1, len(paragraphs) + 1),
                        "text": paragraphs,
                    }
                ),
            )
        )

    for i, table in enumerate(doc.tables, 1):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        if rows:
            width = max(len(r) for r in rows)
            normalized = [
                r + [""] * (width - len(r))
                for r in rows
            ]
            header = normalized[0]
            # Make duplicate/blank headers safe.
            header = [
                str(x).strip() or f"column_{j+1}"
                for j, x in enumerate(header)
            ]
            result.append(
                (
                    f"table_{i}",
                    pd.DataFrame(normalized[1:], columns=header),
                )
            )

    if not result:
        result.append(("data", pd.DataFrame({"text": [""]})))

    return result


def read_pptx(path):
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("PPTX support requires: pip install python-pptx")

    prs = Presentation(path)
    rows = []

    for slide_no, slide in enumerate(prs.slides, 1):
        for shape_no, shape in enumerate(slide.shapes, 1):
            if hasattr(shape, "text") and shape.text.strip():
                rows.append(
                    {
                        "slide": slide_no,
                        "shape": shape_no,
                        "text": shape.text,
                    }
                )

            if getattr(shape, "has_table", False):
                table = shape.table
                for r, row in enumerate(table.rows, 1):
                    rows.append(
                        {
                            "slide": slide_no,
                            "shape": shape_no,
                            "table_row": r,
                            "text": " | ".join(
                                cell.text for cell in row.cells
                            ),
                        }
                    )

    return [
        (
            "content",
            pd.DataFrame(rows)
            if rows
            else pd.DataFrame({"text": []}),
        )
    ]


def read_pdf(path):
    try:
        from pypdf import PdfReader
    except ImportError:
        raise RuntimeError("PDF support requires: pip install pypdf")

    reader = PdfReader(str(path))
    rows = []

    for page_no, page in enumerate(reader.pages, 1):
        try:
            page_text = page.extract_text() or ""
        except Exception as exc:
            page_text = f"[PDF text extraction failed: {exc}]"

        rows.append(
            {
                "page": page_no,
                "text": page_text,
            }
        )

    return [("pages", pd.DataFrame(rows))]


def read_numpy(path):
    import numpy as np

    suffix = path.suffix.lower()

    if suffix == ".npy":
        arr = np.load(path, allow_pickle=False)
        arr = np.asarray(arr)

        if arr.ndim == 0:
            return [("data", pd.DataFrame({"value": [arr.item()]}))]
        if arr.ndim == 1:
            return [("data", pd.DataFrame({"value": arr}))]

        # Flatten all dimensions except the first so it remains tabular.
        return [
            (
                "data",
                pd.DataFrame(
                    arr.reshape(arr.shape[0], -1)
                ),
            )
        ]

    data = np.load(path, allow_pickle=False)
    result = []

    for key in data.files:
        arr = np.asarray(data[key])
        if arr.ndim == 0:
            df = pd.DataFrame({"value": [arr.item()]})
        elif arr.ndim == 1:
            df = pd.DataFrame({"value": arr})
        else:
            df = pd.DataFrame(
                arr.reshape(arr.shape[0], -1)
            )
        result.append((safe_name(key), df))

    return result


def read_mat(path):
    try:
        from scipy.io import loadmat
    except ImportError:
        raise RuntimeError("MAT support requires: pip install scipy")

    obj = loadmat(path)
    result = []

    for key, value in obj.items():
        if key.startswith("__"):
            continue

        try:
            arr = value
            if getattr(arr, "ndim", 0) == 0:
                df = pd.DataFrame({"value": [arr.item()]})
            elif getattr(arr, "ndim", 0) == 1:
                df = pd.DataFrame({"value": arr})
            else:
                df = pd.DataFrame(arr)
            result.append((safe_name(key), df))
        except Exception:
            pass

    if not result:
        result.append(("data", pd.DataFrame({"message": ["No tabular MAT variables found."]})))

    return result


def read_hdf(path):
    try:
        store = pd.HDFStore(path, mode="r")
    except Exception as exc:
        raise RuntimeError(f"Could not open HDF5 file: {exc}")

    result = []
    try:
        for key in store.keys():
            df = store[key]
            result.append((safe_name(key.strip("/")), df))
    finally:
        store.close()

    return result or [("data", pd.DataFrame())]


def read_arff(path):
    try:
        from scipy.io import arff
    except ImportError:
        raise RuntimeError("ARFF support requires: pip install scipy")

    data, _meta = arff.loadarff(path)
    df = pd.DataFrame(data)

    # scipy may return bytes for nominal strings.
    for col in df.columns:
        if df[col].dtype == object:
            df[col] = df[col].map(
                lambda x: x.decode("utf-8", errors="replace")
                if isinstance(x, (bytes, bytearray))
                else x
            )
    return df


def binary_metadata_dataframe(path):
    """
    Universal final fallback for ANY binary/unknown extension.
    It makes the file exportable instead of failing on the extension.
    """
    size = path.stat().st_size if path.exists() else 0

    sha256 = hashlib.sha256()
    preview = b""

    try:
        with open(path, "rb") as f:
            preview = f.read(64 * 1024)
            sha256.update(preview)
            while True:
                chunk = f.read(1024 * 1024)
                if not chunk:
                    break
                sha256.update(chunk)
    except Exception:
        pass

    mime, _ = mimetypes.guess_type(str(path))

    metadata = pd.DataFrame(
        {
            "property": [
                "file_name",
                "extension",
                "mime_type_guess",
                "size_bytes",
                "size_MB",
                "sha256",
            ],
            "value": [
                path.name,
                path.suffix.lower() or "(none)",
                mime or "unknown",
                size,
                round(size / (1024 * 1024), 4),
                sha256.hexdigest(),
            ],
        }
    )

    # A readable byte preview is useful when the unknown file is a
    # proprietary/binary format.
    preview_rows = []
    chunk_size = 16 * 1024
    for offset in range(0, len(preview), chunk_size):
        chunk = preview[offset:offset + chunk_size]
        preview_rows.append(
            {
                "offset": offset,
                "hex_preview": chunk.hex(),
            }
        )

    return [
        ("metadata", metadata),
        ("binary_preview", pd.DataFrame(preview_rows)),
    ]


def decompress_to_temp(path):
    """
    Decompress gzip/bzip2/xz regardless of filename extension.
    Returns a temporary path.
    """
    if looks_like_gzip(path):
        opener = gzip.open
        suffix = ".decompressed"
    elif looks_like_bzip2(path):
        opener = bz2.open
        suffix = ".decompressed"
    elif looks_like_xz(path):
        opener = lzma.open
        suffix = ".decompressed"
    else:
        return None

    temp = Path(
        tempfile.mkstemp(
            prefix="dataset_decompressed_",
            suffix=suffix,
        )[1]
    )

    try:
        with opener(path, "rb") as src, open(temp, "wb") as dst:
            shutil.copyfileobj(src, dst)
        return temp
    except Exception:
        try:
            temp.unlink()
        except Exception:
            pass
        return None


def load_file(path):
    """
    UNIVERSAL loader.

    The file extension is NOT trusted. Detection order:
      1) HTML protection
      2) Office/database/document signatures
      3) archive/compression signatures
      4) well-known structured formats
      5) Markdown/HTML/XML/YAML/JSON
      6) delimited text
      7) plain text
      8) binary metadata fallback

    Therefore an unknown extension such as .bin, .data, .abc, or a URL
    with no extension can still be converted to Excel.
    """
    path = Path(path)

    if is_html_file(path):
        raise ValueError(
            "Downloaded file is HTML, not dataset data."
        )

    name = path.name.lower()
    suffix = path.suffix.lower()

    # --------------------------------------------------------
    # Compressed streams detected by CONTENT, not extension.
    # --------------------------------------------------------
    decompressed = decompress_to_temp(path)
    if decompressed is not None:
        try:
            return load_file(decompressed)
        finally:
            try:
                decompressed.unlink()
            except Exception:
                pass

    # --------------------------------------------------------
    # Office documents before generic ZIP detection.
    # --------------------------------------------------------
    if suffix == ".docx":
        return read_docx(path)

    if suffix == ".pptx":
        return read_pptx(path)

    if suffix == ".pdf" or looks_like_pdf(path):
        return read_pdf(path)

    # --------------------------------------------------------
    # SQLite/database by signature.
    # --------------------------------------------------------
    if suffix in {".db", ".sqlite", ".sqlite3", ".db3"} or looks_like_sqlite(path):
        try:
            return read_sqlite(path)
        except Exception:
            # If it merely has a misleading DB extension, continue sniffing.
            pass

    # --------------------------------------------------------
    # Excel family.
    # --------------------------------------------------------
    if suffix in {".xlsx", ".xls", ".xlsb", ".ods"}:
        return read_excel_universal(path)

    # XLSX/XLSB/ODS are ZIP containers with recognizable signatures.
    if zipfile.is_zipfile(path):
        try:
            with zipfile.ZipFile(path, "r") as z:
                names = set(z.namelist())
            if "[Content_Types].xml" in names and any(
                n.startswith("xl/") for n in names
            ):
                return read_excel_universal(path)
        except Exception:
            pass

    # --------------------------------------------------------
    # Scientific/data formats.
    # --------------------------------------------------------
    if suffix in {".parquet", ".pq"}:
        return [("data", pd.read_parquet(path))]

    if suffix == ".feather":
        return [("data", pd.read_feather(path))]

    if suffix == ".orc":
        return [("data", pd.read_orc(path))]

    if suffix in {".h5", ".hdf5", ".hdf"}:
        return read_hdf(path)

    if suffix in {".dta"}:
        return [("data", pd.read_stata(path))]

    if suffix in {".sav"}:
        return [("data", pd.read_spss(path))]

    if suffix == ".arff":
        return [("data", read_arff(path))]

    if suffix in {".npy", ".npz"}:
        return read_numpy(path)

    if suffix == ".mat":
        return read_mat(path)

    # --------------------------------------------------------
    # Archives by CONTENT, even with a wrong/missing extension.
    # --------------------------------------------------------
    if zipfile.is_zipfile(path) or tarfile.is_tarfile(path):
        return read_archive(path)

    # --------------------------------------------------------
    # Explicit text/structured hints.
    # --------------------------------------------------------
    if name.endswith(".jsonl.gz") or suffix in {".json", ".jsonl", ".ndjson"}:
        return [("data", read_json(path))]

    if suffix in {".yaml", ".yml"}:
        return [("data", read_yaml(path))]

    if suffix == ".xml":
        return [("data", read_xml(path))]

    if suffix in {".md", ".markdown"}:
        try:
            return [("data", read_markdown(path))]
        except Exception:
            # Documentation is still exportable as line-by-line text.
            return [("text", read_plain_text_as_dataframe(path))]

    if suffix in {".html", ".htm"}:
        try:
            tables = pd.read_html(path)
            if tables:
                return [
                    (f"table_{i+1}", df)
                    for i, df in enumerate(tables)
                ]
        except Exception:
            pass
        return [("text", read_plain_text_as_dataframe(path))]

    # Content-based JSON detection, regardless of extension.
    sample = read_text_sample(path, size=200_000)
    stripped = sample.lstrip()

    if stripped.startswith(("{", "[")):
        try:
            return [("data", read_json(path))]
        except Exception:
            pass

    # Content-based XML detection.
    if stripped.startswith("<?xml") or (
        stripped.startswith("<")
        and ">" in stripped[:5000]
    ):
        try:
            return [("data", read_xml(path))]
        except Exception:
            pass

    # Content-based Markdown table detection.
    try:
        md_df = markdown_table_to_dataframe(sample.splitlines())
        if md_df is not None:
            return [("data", md_df)]
    except Exception:
        pass

    # --------------------------------------------------------
    # Generic text -> tabular -> plain-line fallback.
    # --------------------------------------------------------
    if looks_like_text(path):
        try:
            return [("data", read_delimited_universal(path))]
        except Exception:
            return [("text", read_plain_text_as_dataframe(path))]

    # --------------------------------------------------------
    # ANYTHING ELSE: never fail solely because of extension.
    # --------------------------------------------------------
    return binary_metadata_dataframe(path)


# ============================================================
# ARCHIVES
# ============================================================

SUPPORTED = (
    ".csv",
    ".tsv",
    ".txt",
    ".json",
    ".jsonl",
    ".parquet",
    ".xlsx",
    ".xls",
    ".html",
    ".htm",
    ".md",
)


def read_zip(path):
    temp_dir = Path(
        tempfile.mkdtemp(
            prefix="dataset_zip_"
        )
    )

    result = []

    try:
        with zipfile.ZipFile(path, "r") as z:
            bad = z.testzip()

            if bad:
                raise ValueError(
                    f"Corrupt ZIP file: {bad}"
                )

            z.extractall(temp_dir)

        files = [
            p
            for p in temp_dir.rglob("*")
            if (
                p.is_file()
                and not p.name.startswith(".")
                and "__MACOSX" not in p.parts
            )
        ]

        if not files:
            raise ValueError(
                "No supported dataset files were found inside ZIP."
            )

        for file in files:
            try:
                loaded = load_file(file)

                for name, df in loaded:
                    result.append(
                        (
                            safe_name(
                                f"{file.stem}_{name}"
                            ),
                            df
                        )
                    )

            except Exception as exc:
                print(
                    f"Skipping {file}: {exc}"
                )

        return result

    finally:
        shutil.rmtree(
            temp_dir,
            ignore_errors=True
        )


def read_tar(path):
    temp_dir = Path(
        tempfile.mkdtemp(
            prefix="dataset_tar_"
        )
    )

    result = []

    try:
        with tarfile.open(path, "r:*") as tar:
            # Prevent path traversal when extracting untrusted archives.
            root = temp_dir.resolve()

            for member in tar.getmembers():
                target = (temp_dir / member.name).resolve()

                if not str(target).startswith(str(root)):
                    raise ValueError(
                        "Unsafe archive path detected: "
                        f"{member.name}"
                    )

            tar.extractall(temp_dir)

        files = [
            p
            for p in temp_dir.rglob("*")
            if (
                p.is_file()
                and not p.name.startswith(".")
                and "__MACOSX" not in p.parts
            )
        ]

        if not files:
            raise ValueError(
                "No supported dataset files were found inside TAR archive."
            )

        for file in files:
            try:
                loaded = load_file(file)

                for name, df in loaded:
                    result.append(
                        (
                            safe_name(
                                f"{file.stem}_{name}"
                            ),
                            df
                        )
                    )

            except Exception as exc:
                print(
                    f"Skipping {file}: {exc}"
                )

        return result

    finally:
        shutil.rmtree(
            temp_dir,
            ignore_errors=True
        )


def read_archive(path):
    name = path.name.lower()

    if name.endswith(".zip"):
        return read_zip(path)

    if name.endswith(
        (".tar", ".tar.gz", ".tgz")
    ):
        return read_tar(path)

    raise ValueError(
        f"Unsupported archive format: {path}"
    )


# ============================================================
# CLEAN DATAFRAME
# ============================================================

def clean_dataframe(df):

    df = df.copy()

    # Duplicate columns
    columns = []
    counts = {}

    for col in df.columns:

        col = str(col)

        counts[col] = (
            counts.get(col, 0) + 1
        )

        if counts[col] == 1:
            columns.append(col)
        else:
            columns.append(
                f"{col}_{counts[col]}"
            )

    df.columns = columns

    # Remove control characters from text
    for col in df.columns:

        if df[col].dtype == "object":

            df[col] = df[col].map(
                lambda value:
                    re.sub(
                        r"[\x00-\x08\x0B\x0C\x0E-\x1F]",
                        "",
                        value
                    )
                    if isinstance(
                        value,
                        str
                    )
                    else value
            )

    return df


# ============================================================
# SAVE EXCEL / CSV / PARQUET
# ============================================================

def save_dataframe(
    df,
    base_name,
    save_extra=True
):

    df = clean_dataframe(df)

    base_name = safe_name(
        base_name
    )

    # -------------------------
    # Excel
    # -------------------------

    excel_path = unique_path(
        OUTPUT_DIR
        / f"{base_name}.xlsx"
    )

    print(
        f"\nSaving Excel:"
        f"\n{excel_path}"
    )

    print(
        f"Rows: {len(df):,}"
    )

    print(
        f"Columns: {list(df.columns)}"
    )

    with pd.ExcelWriter(
        excel_path,
        engine="openpyxl"
    ) as writer:

        if len(df) == 0:

            df.to_excel(
                writer,
                sheet_name="data",
                index=False
            )

        else:

            part = 1

            for start in range(
                0,
                len(df),
                EXCEL_MAX_ROWS
            ):

                end = min(
                    start + EXCEL_MAX_ROWS,
                    len(df)
                )

                chunk = df.iloc[
                    start:end
                ]

                sheet_name = (
                    "data"
                    if part == 1
                    else f"data_{part}"
                )

                chunk.to_excel(
                    writer,
                    sheet_name=sheet_name[:31],
                    index=False
                )

                print(
                    f"Excel sheet {sheet_name}: "
                    f"rows {start+1:,} - {end:,}"
                )

                part += 1

    print(
        "SUCCESS Excel:",
        excel_path
    )

    # -------------------------
    # CSV
    # -------------------------

    if save_extra:

        csv_path = unique_path(
            OUTPUT_DIR
            / f"{base_name}.csv"
        )

        print(
            "Saving CSV:",
            csv_path
        )

        df.to_csv(
            csv_path,
            index=False,
            encoding="utf-8-sig"
        )

        print(
            "SUCCESS CSV:",
            csv_path
        )

        # -------------------------
        # Parquet
        # -------------------------

        parquet_path = unique_path(
            OUTPUT_DIR
            / f"{base_name}.parquet"
        )

        try:

            df.to_parquet(
                parquet_path,
                index=False
            )

            print(
                "SUCCESS Parquet:",
                parquet_path
            )

        except Exception as exc:

            print(
                "Parquet skipped:",
                exc
            )


# ============================================================
# HUGGING FACE
# ============================================================

def get_hf_token():

    return (
        os.getenv("HF_TOKEN")
        or os.getenv(
            "HUGGINGFACE_TOKEN"
        )
    )


def load_huggingface_dataset(
    dataset_id
):

    try:

        from datasets import (
            load_dataset,
            get_dataset_config_names,
        )

    except ImportError:

        raise RuntimeError(
            "Install Hugging Face packages:\n"
            "pip install datasets huggingface_hub"
        )

    token = get_hf_token()

    print(
        "\n"
        + "=" * 75
    )

    print(
        "HUGGING FACE DATASET"
    )

    print(
        "=" * 75
    )

    print(
        "Dataset:",
        dataset_id
    )

    # Configs
    try:

        configs = (
            get_dataset_config_names(
                dataset_id,
                token=token
            )
        )

    except Exception as exc:

        print(
            "Could not get configs:",
            exc
        )

        configs = [
            "default"
        ]

    if not configs:
        configs = ["default"]

    print(
        "\nAvailable configs:"
    )

    for i, config in enumerate(
        configs,
        1
    ):

        print(
            f"  {i}. {config}"
        )

    if len(configs) > 1:
        print(
            "  A. ALL configs"
        )

    if len(configs) == 1:

        selected_configs = configs

    else:

        choice = input(
            "\nChoose config number "
            "or A for all: "
        ).strip()

        if choice.lower() == "a":

            selected_configs = configs

        else:

            try:

                selected_configs = [
                    configs[
                        int(choice) - 1
                    ]
                ]

            except Exception:

                print(
                    "Invalid choice. "
                    "Using all configs."
                )

                selected_configs = configs

    for config in selected_configs:

        print(
            "\n"
            + "-" * 75
        )

        print(
            "Config:",
            config
        )

        print(
            "-" * 75
        )

        try:

            if config == "default":

                ds = load_dataset(
                    dataset_id,
                    token=token
                )

            else:

                ds = load_dataset(
                    dataset_id,
                    name=config,
                    token=token
                )

        except Exception as exc:

            print(
                "Could not load config:",
                exc
            )

            continue

        splits = list(
            ds.keys()
        )

        print(
            "\nAvailable splits:"
        )

        for i, split in enumerate(
            splits,
            1
        ):

            print(
                f"  {i}. {split}"
            )

        print(
            "  A. ALL splits"
        )

        if len(splits) == 1:

            selected_splits = splits

        else:

            choice = input(
                "\nChoose split number "
                "or A for all: "
            ).strip()

            if choice.lower() == "a":

                selected_splits = splits

            else:

                try:

                    selected_splits = [
                        splits[
                            int(choice) - 1
                        ]
                    ]

                except Exception:

                    selected_splits = splits

        for split in selected_splits:

            print(
                "\nLoading split:",
                split
            )

            try:

                df = (
                    ds[split]
                    .to_pandas()
                )

                print(
                    f"Rows: {len(df):,}"
                )

                print(
                    "Columns:",
                    list(df.columns)
                )

                dataset_name = safe_name(
                    dataset_id.split("/")[-1]
                )

                config_name = safe_name(
                    config
                )

                split_name = safe_name(
                    split
                )

                base_name = (
                    f"{dataset_name}_"
                    f"{config_name}_"
                    f"{split_name}"
                )

                save_dataframe(
                    df,
                    base_name,
                    save_extra=SAVE_EXTRA
                )

            except Exception as exc:

                print(
                    f"ERROR in split "
                    f"{split}: {exc}"
                )


# ============================================================
# URL PROCESSING
# ============================================================

def process_url(url):

    url = url.strip()

    # --------------------------------------------------------
    # GitHub documentation (README/data.md/etc.)
    # --------------------------------------------------------

    if is_github_url(url):
        path_lower = urlparse(url).path.lower()

        if path_lower.endswith(
            (".md", ".markdown", "/readme")
        ):
            dataset_urls = find_github_dataset_urls(url)

            if dataset_urls:
                print(
                    "\nDataset URL(s) discovered from "
                    "GitHub documentation:"
                )

                for i, dataset_url in enumerate(
                    dataset_urls,
                    1
                ):
                    print(
                        f"  {i}. {dataset_url}"
                    )

                for dataset_url in dataset_urls:
                    process_url(dataset_url)

                return

            # If it is a Markdown table, download and parse it.
            raw_url = github_to_raw(url)
            path = download_http(raw_url)

            dataframes = load_file(path)

            for name, df in dataframes:
                base_name = safe_name(
                    f"{dataset_name_from_url(url)}_{name}"
                )

                save_dataframe(
                    df,
                    base_name,
                    save_extra=SAVE_EXTRA
                )

            return

    # --------------------------------------------------------
    # GitHub normal/blob URL
    # --------------------------------------------------------

    if is_github_url(url):

        raw_url = github_to_raw(
            url
        )

        if raw_url != url:

            print(
                "\n"
                "GitHub BLOB URL detected."
            )

            print(
                "Original:"
            )

            print(url)

            print(
                "\nConverted automatically "
                "to RAW dataset URL:"
            )

            print(raw_url)

            url = raw_url

    # --------------------------------------------------------
    # Hugging Face dataset page
    # --------------------------------------------------------

    if (
        "huggingface.co"
        in urlparse(url).netloc.lower()
        and "/datasets/" in urlparse(url).path
        and "/resolve/" not in urlparse(url).path
        and "/raw/" not in urlparse(url).path
    ):

        dataset_id = (
            hf_dataset_url_to_id(url)
        )

        if dataset_id:

            load_huggingface_dataset(
                dataset_id
            )

            return

    # --------------------------------------------------------
    # Download
    # --------------------------------------------------------

    path = download_http(
        url
    )

    # --------------------------------------------------------
    # Archives
    # --------------------------------------------------------

    if is_archive(path) and not is_office_zip(path):

        dataframes = read_archive(
            path
        )

    else:

        dataframes = load_file(
            path
        )

    if not dataframes:

        raise RuntimeError(
            "No dataset table found."
        )

    base_name = (
        dataset_name_from_url(
            url
        )
    )

    for name, df in dataframes:

        final_name = safe_name(
            f"{base_name}_{name}"
        )

        save_dataframe(
            df,
            final_name,
            save_extra=SAVE_EXTRA
        )


# ============================================================
# LOCAL FILE
# ============================================================

def process_local_file(path):

    print(
        "\nLocal file:",
        path
    )

    if is_archive(path) and not is_office_zip(path):

        dataframes = read_archive(
            path
        )

    else:

        dataframes = load_file(
            path
        )

    base_name = safe_name(
        path.stem
    )

    for name, df in dataframes:

        save_dataframe(
            df,
            safe_name(
                f"{base_name}_{name}"
            ),
            save_extra=SAVE_EXTRA
        )


# ============================================================
# MAIN
# ============================================================

def main():

    global SAVE_EXTRA

    print(
        "\n"
        + "=" * 75
    )

    print(
        " UNIVERSAL DATASET DOWNLOADER v5"
    )

    print(
        " GitHub + Hugging Face + CSV + TSV + JSON + Parquet"
    )

    print(
        " Excel + archives + documents + unknown extensions"
    )

    print(
        "=" * 75
    )

    print(
        "\nOutput folder:"
    )

    print(
        OUTPUT_DIR.resolve()
    )

    print(
        "\nUniversal format mode: ON"
    )

    print(
        "The extension is only a hint; content/signature detection is used."
    )

    print(
        "\nExamples:"
    )

    print(
        "\n1. Hugging Face ID:"
    )

    print(
        "   qiaojin/PubMedQA"
    )

    print(
        "\n2. Hugging Face URL:"
    )

    print(
        "   https://huggingface.co/datasets/qiaojin/PubMedQA"
    )

    print(
        "\n3. GitHub BLOB URL:"
    )

    print(
        "   https://github.com/user/repo/blob/main/data/file.tsv"
    )

    print(
        "\n4. Direct file URL:"
    )

    print(
        "   https://example.com/data.csv"
    )

    print(
        "\n5. Local file (any extension):"
    )

    print(
        r"   D:\data\anything.xyz"
    )

    print(
        "\n6. URL with unknown extension:"
    )

    print(
        "   https://example.com/download?id=123"
    )

    print(
        "\n7. SQLite database:"
    )

    print(
        r"   D:\data\database.db"
    )

    value = input(
        "\nDataset ID / URL / local file path: "
    ).strip().strip('"').strip("'")

    if not value:

        print(
            "Nothing entered."
        )

        return

    extra = input(
        "\nAlso save CSV + Parquet? [Y/n]: "
    ).strip().lower()

    SAVE_EXTRA = (
        extra
        not in ("n", "no")
    )

    print(
        "\n"
        + "=" * 75
    )

    try:

        # Local file
        local = Path(value)

        if (
            local.exists()
            and local.is_file()
        ):

            process_local_file(
                local
            )

        # Hugging Face dataset ID
        elif looks_like_hf_id(
            value
        ):

            load_huggingface_dataset(
                value
            )

        # URL
        elif is_url(
            value
        ):

            process_url(
                value
            )

        else:

            raise ValueError(
                "Input is not a valid "
                "Hugging Face ID, URL, "
                "or local file."
            )

        print(
            "\n"
            + "=" * 75
        )

        print(
            "DONE"
        )

        print(
            "Files saved in:"
        )

        print(
            OUTPUT_DIR.resolve()
        )

        print(
            "=" * 75
        )

    except KeyboardInterrupt:

        print(
            "\nStopped by user."
        )

    except Exception as exc:

        print(
            "\n"
            + "=" * 75
        )

        print(
            "ERROR"
        )

        print(
            "=" * 75
        )

        print(
            type(exc).__name__,
            ":",
            exc
        )

        print(
            "\nImportant:"
        )

        print(
            "- GitHub blob URLs are automatically "
            "converted to raw URLs."
        )

        print(
            "- GitHub README/data.md documentation can "
            "be inspected for dataset download URLs/scripts."
        )

        print(
            "- ZIP/TAR/TAR.GZ archives are automatically "
            "extracted and processed."
        )

        print(
            "- If a website requires login, "
            "CAPTCHA, API key, or permission, "
            "the downloader cannot bypass it."
        )

        print(
            "- For GitHub, make sure the repository "
            "and file are publicly accessible."
        )

        print(
            "- Unknown extensions are detected by file content; "
            "if decoding is impossible, file metadata/byte preview "
            "is exported instead of failing."
        )

        print(
            "\nOutput folder:"
        )

        print(
            OUTPUT_DIR.resolve()
        )


if __name__ == "__main__":
    main()